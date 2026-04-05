#!/usr/bin/env python3
"""Generate a board-level PowerPoint for the Real-Time Testing Dashboard (requires python-pptx)."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def _title_only(prs: Presentation, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def _bullets(prs: Presentation, title: str, lines: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(lines):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18 if i == 0 else 16)
    if body.paragraphs:
        body.paragraphs[0].font.size = Pt(16)


def _two_column_bullets(prs: Presentation, title: str, left: list[str], right: list[str]) -> None:
    """Title + two text boxes (manual shapes) for comparison slides."""
    blank = prs.slide_layouts[6]  # often blank
    slide = prs.slides.add_slide(blank)
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0x15, 0x1D, 0x35)

    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.5))
    ltf = left_box.text_frame
    ltf.clear()
    for i, line in enumerate(left):
        p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.level = 0

    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(5.5))
    rtf = right_box.text_frame
    rtf.clear()
    for i, line in enumerate(right):
        p = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.level = 0


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _title_only(
        prs,
        "Real-Time Testing Dashboard",
        "QA observability for live test execution & CI health\nBoard briefing",
    )

    _bullets(
        prs,
        "Executive summary",
        [
            "Single pane of glass for automated test execution: pass rates, defects, and trends—updated continuously.",
            "Connects CI (e.g. GitHub Actions + Playwright) to a live dashboard without manual report chasing.",
            "Modern stack (FastAPI + React), cloud-ready (e.g. Render + Vercel), API-first for future integrations.",
            "Designed to shorten feedback loops and improve visibility for engineering and quality leadership.",
        ],
    )

    _bullets(
        prs,
        "The problem we solve",
        [
            "Test results trapped in CI logs, artifacts, and siloed tools—hard to see at a glance.",
            "Leadership lacks a real-time view of quality across environments and suites.",
            "Playwright/HTML reports are not linkable from GitHub artifacts without extra steps.",
            "Teams need a lightweight, defensible view of quality without heavy enterprise tooling.",
        ],
    )

    _bullets(
        prs,
        "Solution: Real-Time Testing Dashboard",
        [
            "Live KPIs: total runs, cases, pass rate, open defects, execution status mix.",
            "Trend views: module quality, environment distribution, latest run feed.",
            "Real-time updates via WebSocket; optional CI ingestion for production runs.",
            "Embedded HTML reports when CI uploads the report bundle to the API—full traceability in one place.",
        ],
    )

    _two_column_bullets(
        prs,
        "Business outcomes",
        [
            "Faster triage",
            "• See failures and regressions sooner",
            "",
            "Better transparency",
            "• Shared view for Eng & QA leadership",
            "",
            "Lower friction",
            "• Open patterns; fits existing CI",
        ],
        [
            "Operational clarity",
            "• Environment and module trends",
            "",
            "Audit-friendly",
            "• Runs and cases stored in Postgres",
            "",
            "Cost control",
            "• Fits free/low-cost cloud tiers",
        ],
    )

    _bullets(
        prs,
        "Product capabilities (current)",
        [
            "REST API: summary, runs, case updates; WebSocket channel for live summary refresh.",
            "GitHub Actions ingestion: signed POST of test results + optional Playwright HTML report zip.",
            "Dashboard UI: KPI cards, execution & environment charts, module quality, live run table, HTML embed.",
            "Demo mode: synthetic runs for stakeholder demos without production data.",
        ],
    )

    _bullets(
        prs,
        "Architecture (high level)",
        [
            "Frontend: React + TypeScript (Vite)—static or CDN-hosted; talks to API over HTTPS/WSS.",
            "Backend: FastAPI, SQLAlchemy, Alembic migrations; Postgres in production.",
            "Ingestion: secure token (X-Ingest-Token); CI never exposes secrets to the browser.",
            "Horizontal path: CI → API → database → dashboard + WebSocket broadcast to viewers.",
        ],
    )

    _bullets(
        prs,
        "Security & governance",
        [
            "Ingest tokens and credentials live only on the server (e.g. Render secrets), not in the frontend.",
            "CORS restricted to known dashboard origins in production.",
            "API-driven model: role-based access can be layered via gateway or future auth without UI rewrite.",
        ],
    )

    _bullets(
        prs,
        "Roadmap & investment themes (illustrative)",
        [
            "Deeper CI integrations (additional reporters, richer metadata).",
            "Optional auth/RBAC and team-scoped projects for multi-tenant use.",
            "Alerting (Slack/email) on pass-rate or defect thresholds.",
            "Longer retention policies and export for compliance reporting.",
        ],
    )

    _bullets(
        prs,
        "Success metrics (how we measure value)",
        [
            "Time to detect a regression suite failure (median).",
            "Adoption: active CI workflows posting to the dashboard.",
            "Reduction in ad-hoc “where is the report?” escalations.",
            "Leadership usage: weekly active viewers of the live dashboard.",
        ],
    )

    _bullets(
        prs,
        "Why now",
        [
            "Quality is a board-level topic: velocity must not outrun visibility.",
            "A lightweight, open, API-first dashboard aligns with modern DevOps and FinOps discipline.",
            "Ready to pilot with one team, expand as ingestion patterns stabilize.",
        ],
    )

    _title_only(prs, "Thank you", "Questions & discussion")

    return prs


def main() -> None:
    out_dir = Path(__file__).resolve().parent.parent / "presentations"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "realtime-testing-dashboard-board.pptx"
    prs = build()
    prs.save(out_path)
    print(f"Wrote {out_path}")


if __name__ == "__main__":
    main()
